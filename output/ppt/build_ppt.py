from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path

base = Path('/home/panzek/project-menuju-sukses/inventori-pergudangan')
img1 = base / 'output/playwright/warehouse-slide-shot.png'
img2 = base / 'output/playwright/petayuai-slide-shot.png'
out = base / 'output/ppt/materi-manajemen-gudang-petayuai-2-slide.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 28, 45)


def add_bullets(slide, lines):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(5.6), Inches(5.9))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)


def add_screenshot(slide, path):
    # fit screenshot on right side
    slide.shapes.add_picture(str(path), Inches(6.2), Inches(0.95), width=Inches(6.55), height=Inches(5.9))

# Slide 1
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s1, 'Fitur Manajemen Gudang (/warehouse) - Poin Inti')
add_bullets(s1, [
    'Kelola master zona dan rak: tambah, ubah, hapus, kapasitas, status.',
    'Kelola stok per rak untuk akurasi lokasi barang dan proses picking.',
    'Dashboard real-time: total zona, rak, item, kapasitas, dan okupansi.',
    'Layout editor: atur posisi/ukuran/rotasi, simpan snapshot, ekspor PDF.',
    'Hak akses: Manager (master & layout), Supervisor (operasional stok rak).',
])
add_screenshot(s1, img1)

# Slide 2
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s2, 'PETAYU AI + Dampak Operasional')
add_bullets(s2, [
    'PETAYU AI memberi insight operasional gudang secara real-time.',
    'Use case cepat: ringkasan stok, produk stok rendah, prediksi stok 3-7 hari.',
    'Membantu prioritas tindakan: restock, redistribusi rak, evaluasi efisiensi.',
    'Output AI jadi early insight, eksekusi tetap di modul warehouse/inventory.',
    'Dampak: keputusan lebih cepat, utilisasi ruang lebih terukur, audit lebih rapi.',
])
add_screenshot(s2, img2)

prs.save(out)
print(out)
